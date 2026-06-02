import io
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

def _set_cell_border(cell, color="000000", width="12700"):
    """PPT 표 셀에 검정색 테두리를 적용하기 위해 XML을 직접 제어합니다."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    
    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln_existing = tcPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{border_name}')
        if ln_existing is not None:
            tcPr.remove(ln_existing)
            
        ln = OxmlElement(f'a:{border_name}')
        ln.set('w', width)
        ln.set('cmpd', 'sng')
        
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color)
        
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)

# PPTX에서 사용할 커스텀 색상 매핑
COLOR_RED = RGBColor(255, 0, 0)
COLOR_YELLOW = RGBColor(255, 255, 0)
COLOR_GREEN = RGBColor(0, 176, 80)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_HEADER_BG = RGBColor(230, 230, 230)

def _set_cell_style(cell, text, bg_color=None, text_color=COLOR_BLACK, bold=False, font_size=9, align=PP_ALIGN.CENTER):
    """PPT 표의 셀 스타일을 적용합니다."""
    cell.text = text
    
    # 텍스트 포맷팅
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = text_color
            run.font.bold = bold
            
    # 배경색 적용
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        
def _get_fill_color_for_score(earned_pt, max_pt):
    if max_pt <= 0:
        return COLOR_YELLOW
    ratio = earned_pt / max_pt
    if ratio >= 0.8:
        return COLOR_GREEN
    elif ratio >= 0.6:
        return COLOR_YELLOW
    else:
        return COLOR_RED

def _fmt_score(v):
    return f"{v:.1f}점" if v > 0 else "-"
    
def _fmt_num(v):
    return f"{v:,.1f}"

def export_p1_summary_ppt_bytes(data: dict, cat_filter: str, metrics_filter=None):
    """
    P1 대시보드의 지점 요약본을 PPT(python-pptx) 형태로 추출합니다.
    report_generator.py의 export_p1_summary_excel_bytes 구조를 참조하여
    데이터를 집계한 뒤 PPT 표로 렌더링합니다.
    """
    if metrics_filter is None:
        metrics_filter = ["할인율", "BEST상품", "신선도", "시즌", "아이템"]
        
    stores = data.get("STORES", [])
    brands = data.get("BRANDS", [])
    detail = data.get("DETAIL", {})
    
    if not stores or not brands:
        return None

    metrics = ["dis", "best", "fresh", "season", "item"]
    agg_brands = []
    agg_detail = {}

    for store in stores:
        store_brands = [b for b in brands if b.get("store") == store and b.get("category") == cat_filter]
        if not store_brands:
            continue

        b_count = len(store_brands)
        agg_tM = sum(float(b.get("tM", 0.0)) for b in store_brands)
        agg_sM = sum(float(b.get("sM", 0.0)) for b in store_brands)
        
        rep_b = store_brands[0]
        agg_b = dict(rep_b)
        agg_b["name"] = "지점 합계"
        agg_b["type_label"] = "전체"
        agg_b["sM"] = agg_sM
        agg_b["tM"] = agg_tM
        agg_b["is_p1_aggregate"] = True

        for m in metrics:
            if agg_tM > 0:
                agg_m_val = sum(float(b.get(m, 0.0)) * (float(b.get("tM", 0.0)) / agg_tM) for b in store_brands)
            else:
                agg_m_val = sum(float(b.get(m, 0.0)) for b in store_brands) / b_count if b_count > 0 else 0.0
            agg_b[m] = agg_m_val

        if store not in agg_detail:
            agg_detail[store] = {}
        agg_detail[store]["지점 합계"] = {"전체": {}}
        store_det = agg_detail[store]["지점 합계"]["전체"]

        agg_calc = 0.0
        agg_max_weights = 0.0
        for m in metrics:
            store_det[m] = {"segs": []}
            seg_map = {}
            sum_weights = sum(float(b.get("tM", 0.0)) for b in store_brands)
            
            for b in store_brands:
                b_name = b.get("name")
                b_type = b.get("type_label")
                b_tM = float(b.get("tM", 0.0))
                b_weight = (b_tM / sum_weights) if sum_weights > 0 else (1.0 / len(store_brands))
                
                s_weights = b.get("scoring_guide", {}).get("score_weights", {})
                w_key = "sea" if m == "season" else m
                m_weight = s_weights.get(w_key, 0.0)
                
                m_data = detail.get(store, {}).get(b_name, {}).get(b_type, {}).get(m, {})
                segs_raw = m_data.get("segs", [])
                
                seg_opt_pct_map = {}
                sum_opt_pct = 0.0
                for s in segs_raw:
                    if isinstance(s, dict) and "key" in s:
                        op_val = float(s.get("opt_pct", 0.0))
                        seg_opt_pct_map[s["key"]] = op_val
                        sum_opt_pct += op_val
                if sum_opt_pct <= 0:
                    sum_opt_pct = 100.0
                    
                for s in segs_raw:
                    if not isinstance(s, dict) or "key" not in s:
                        continue
                    k = s["key"]
                    brand_opt_pct = seg_opt_pct_map.get(k, 0.0)
                    brand_seg_max_pt = m_weight * (brand_opt_pct / sum_opt_pct)
                    
                    valM = float(s.get("valM", 0.0))
                    targetM = float(s.get("targetM", 0.0))
                    pct = float(s.get("pct", 0.0))
                    
                    real_pct = (valM / targetM * 100.0) if targetM > 0 else pct
                    is_over_120 = (real_pct > 120.0)
                    
                    capped_pct = min(pct, 100.0)
                    raw_earned_pt = brand_seg_max_pt * (capped_pct / 100.0)
                    earned_pt = raw_earned_pt * 0.9 if is_over_120 else raw_earned_pt
                    earned_pt = min(earned_pt, brand_seg_max_pt)
                    
                    if k not in seg_map:
                        seg_map[k] = {
                            "valM": 0.0,
                            "targetM": 0.0,
                            "opt_pct": float(s.get("opt_pct", 0.0)),
                            "earned_pt": 0.0,
                        }
                    seg_map[k]["valM"] += valM
                    seg_map[k]["targetM"] += targetM
                    seg_map[k]["earned_pt"] += earned_pt * b_weight

            m_total_earned = 0.0
            avg_m_weight = (
                sum(b.get("scoring_guide", {}).get("score_weights", {}).get("sea" if m == "season" else m, 0.0) * (float(b.get("tM", 0.0)) / sum_weights)
                    for b in store_brands) if sum_weights > 0 else
                sum(b.get("scoring_guide", {}).get("score_weights", {}).get("sea" if m == "season" else m, 0.0)
                    for b in store_brands) / len(store_brands)
            )
            for k, v in seg_map.items():
                agg_real_pct = (v["valM"] / v["targetM"] * 100.0) if v["targetM"] > 0 else 0.0
                agg_is_over_120 = (agg_real_pct > 120.0)
                m_total_earned += v["earned_pt"]
                store_det[m]["segs"].append({
                    "key": k,
                    "valM": v["valM"],
                    "targetM": v["targetM"],
                    "opt_pct": v["opt_pct"],
                    "pct": agg_real_pct,
                    "earned_pt": v["earned_pt"],
                    "is_over_120": agg_is_over_120,
                })
            m_earned = min(m_total_earned, avg_m_weight)
            
            metric_filter_map = {"dis": "할인율", "best": "BEST상품", "fresh": "신선도", "season": "시즌", "item": "아이템"}
            is_all = ("전체" in metrics_filter) or (len(metrics_filter) >= 5)
            
            if is_all or metric_filter_map[m] in metrics_filter:
                agg_calc += m_earned
                agg_max_weights += avg_m_weight
        
        if agg_max_weights > 0 and agg_max_weights < 100.0 and not (("전체" in metrics_filter) or (len(metrics_filter) >= 5)):
            normalized_calc = (agg_calc / agg_max_weights) * 100.0
        else:
            normalized_calc = agg_calc
            
        agg_b["calculated_total"] = round(min(normalized_calc, 100.0), 1)
        agg_brands.append(agg_b)

    if not agg_brands:
        return None

    # 데이터 정렬 (총점 내림차순)
    sorted_brands = sorted(agg_brands, key=lambda x: x.get("calculated_total", 0.0), reverse=True)
    
    # ── PPT 객체 생성 ──
    prs = Presentation()
    # A4 가로 규격 세팅 (297mm * 210mm)
    prs.slide_width = Mm(297)
    prs.slide_height = Mm(210)
    
    slide_layout = prs.slide_layouts[5] # Title only
    
    # 표 차원 계산을 위한 metrics_config 구성
    from core.report_generator import _extract_seg_keys_labels
    sample_b = sorted_brands[0] if sorted_brands else {}
    s_store = sample_b.get("store", "")
    sample_detail = agg_detail.get(s_store, {}).get("지점 합계", {}).get("전체", {})
    
    item_seg_keys, item_labels = _extract_seg_keys_labels(sample_detail, "item", ["Outer", "Top", "Bottom", "Skirt", "Dress"], ["아우터", "상의", "하의", "스커트", "원피스"])
    dis_keys, dis_labels = _extract_seg_keys_labels(sample_detail, "dis", ["d70", "d50", "d30", "d10", "d0"], ["70%이상", "50-69%", "30-49%", "1-29%", "정상가"])
    fresh_keys, fresh_labels = _extract_seg_keys_labels(sample_detail, "fresh", ["new", "plan"], ["신상", "기획"])
    season_keys, season_labels = _extract_seg_keys_labels(sample_detail, "season", ["spring", "summer"], ["봄(SS)", "여름(SS)"])

    metrics_config = {}
    if "할인율" in metrics_filter: metrics_config["dis"] = {"title": "할인율", "keys": dis_keys, "labels": dis_labels}
    if "BEST상품" in metrics_filter: metrics_config["best"] = {"title": "BEST상품", "keys": ["best"], "labels": ["판매\nTOP10"]}
    if "신선도" in metrics_filter: metrics_config["fresh"] = {"title": "신선도", "keys": fresh_keys, "labels": fresh_labels}
    if "시즌" in metrics_filter: metrics_config["season"] = {"title": "시즌", "keys": season_keys, "labels": season_labels}
    if "아이템" in metrics_filter: metrics_config["item"] = {"title": "아이템", "keys": item_seg_keys, "labels": item_labels}
    
    for m_id, m_info in metrics_config.items():
        w_key = "sea" if m_id == "season" else m_id
        g_weight = sample_b.get("scoring_guide", {}).get("score_weights", {}).get(w_key, 0.0)
        m_info["max_w"] = g_weight
        num_keys = len(m_info["keys"])
        m_info["seg_max_pts"] = {k: g_weight/num_keys for k in m_info["keys"]} if num_keys > 0 else {}
 
    common_headers = [("랭크", 1), ("복종", 1), ("지점", 1), ("총\n점수\n(100점)", 1), ("총보유\n재고액", 1)]
    total_cols = sum(span for _, span in common_headers) + sum(len(m_info["keys"]) + 1 for m_info in metrics_config.values())
    
    # ── 10개 지점 단위 청킹(Paging) 생성 ──
    chunk_size = 10
    num_chunks = math.ceil(len(sorted_brands) / chunk_size) if sorted_brands else 1
    
    for chunk_idx in range(num_chunks):
        start_idx = chunk_idx * chunk_size
        end_idx = min(start_idx + chunk_size, len(sorted_brands))
        current_chunk = sorted_brands[start_idx:end_idx]
        
        slide = prs.slides.add_slide(slide_layout)
        
        # 타이틀
        title_shape = slide.shapes.title
        title_shape.text = f"[한국유통] 상품구색 노출/측정 (핵심 {len(sorted_brands)}개점 {cat_filter}) - {chunk_idx + 1}/{num_chunks}"
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(24)
                
        rows = 3 + len(current_chunk)
        # A4 기준 균등 여백 설정 (양 옆 5mm, 상단 28mm, 하단 15mm 여백)
        x, y = Mm(5), Mm(28)
        cx, cy = Mm(287), Mm(167)
        
        shape = slide.shapes.add_table(rows, total_cols, x, y, cx, cy)
        table = shape.table
        
        base_col_w = cx / total_cols
        for i in range(total_cols):
            table.columns[i].width = int(base_col_w)

        # 헤더 그리기
        col_idx = 0
        for hdr_title, colspan in common_headers:
            cell = table.cell(0, col_idx)
            _set_cell_style(cell, hdr_title, bg_color=COLOR_HEADER_BG, bold=True)
            cell.merge(table.cell(2, col_idx))
            col_idx += 1
            
        for m_id, m_info in metrics_config.items():
            colspan = len(m_info["keys"]) + 1
            cell_h1 = table.cell(0, col_idx)
            _set_cell_style(cell_h1, f"{m_info['title']}({int(m_info['max_w'])}점)", bg_color=COLOR_HEADER_BG, bold=True)
            if colspan > 1:
                cell_h1.merge(table.cell(0, col_idx + colspan - 1))
                
            cell_h2_sum = table.cell(1, col_idx)
            _set_cell_style(cell_h2_sum, "합계", bg_color=COLOR_HEADER_BG, bold=True)
            
            for i, lbl in enumerate(m_info["labels"]):
                cell_lbl = table.cell(1, col_idx + 1 + i)
                _set_cell_style(cell_lbl, lbl.replace("\n", ""), bg_color=COLOR_HEADER_BG, bold=True, font_size=8)
                
            cell_h3_sum = table.cell(2, col_idx)
            _set_cell_style(cell_h3_sum, f"{int(m_info['max_w'])}점", bg_color=COLOR_HEADER_BG, bold=True)
            for i, k in enumerate(m_info["keys"]):
                seg_w = m_info["seg_max_pts"].get(k, 0)
                cell_h3_seg = table.cell(2, col_idx + 1 + i)
                _set_cell_style(cell_h3_seg, f"{int(seg_w)}점", bg_color=COLOR_HEADER_BG, bold=True)
                
            col_idx += colspan

        # 데이터 그리기
        for r_idx, b in enumerate(current_chunk):
            row_i = 3 + r_idx
            store = b.get("store", "")
            b_name = b.get("name", "")
            b_type = b.get("type_label", "")
            b_detail = agg_detail.get(store, {}).get(b_name, {}).get(b_type, {})
            
            c_idx = 0
            # 랭크
            _set_cell_style(table.cell(row_i, c_idx), str(start_idx + r_idx + 1))
            c_idx += 1
            # 복종
            _set_cell_style(table.cell(row_i, c_idx), b.get("category", ""))
            c_idx += 1
            # 지점
            _set_cell_style(table.cell(row_i, c_idx), store)
            c_idx += 1
            
            # 총점수
            tot_score = b.get("calculated_total", 0.0)
            bg_col = _get_fill_color_for_score(tot_score, 100.0)
            txt_col = COLOR_WHITE if bg_col == COLOR_RED else COLOR_BLACK
            _set_cell_style(table.cell(row_i, c_idx), f"{int(round(tot_score))}점", bg_color=bg_col, text_color=txt_col, bold=True)
            c_idx += 1
            
            # 총재고
            s_val = float(b.get("sM", 0.0))
            _set_cell_style(table.cell(row_i, c_idx), _fmt_num(s_val), bg_color=COLOR_YELLOW)
            c_idx += 1
            
            # 지표별 데이터
            for m_id, m_info in metrics_config.items():
                m_data = b_detail.get(m_id) or {}
                segs_raw = m_data.get("segs", [])
                seg_by_key = {s["key"]: s for s in segs_raw if isinstance(s, dict) and "key" in s}
                
                m_earned = float(b.get(m_id, 0.0))
                g_weight = m_info["max_w"]
                bg_col = _get_fill_color_for_score(m_earned, g_weight)
                txt_col = COLOR_WHITE if bg_col == COLOR_RED else COLOR_BLACK
                _set_cell_style(table.cell(row_i, c_idx), f"{_fmt_score(m_earned)}", bg_color=bg_col, text_color=txt_col, bold=True)
                c_idx += 1
                
                for k in m_info["keys"]:
                    seg = seg_by_key.get(k)
                    seg_max = m_info["seg_max_pts"].get(k, 0)
                    if seg:
                        valM = float(seg.get("valM", 0.0))
                        earned_pt = float(seg.get("earned_pt", 0.0))
                        is_over = seg.get("is_over_120", False)
                        
                        bg_col = _get_fill_color_for_score(earned_pt, seg_max)
                        txt_col = COLOR_WHITE if bg_col == COLOR_RED else COLOR_BLACK
                        if is_over:
                            bg_col = COLOR_YELLOW
                            txt_col = COLOR_BLACK
                            
                        _set_cell_style(table.cell(row_i, c_idx), f"{_fmt_num(valM)}\n({_fmt_score(earned_pt)})", bg_color=bg_col, text_color=txt_col, font_size=8)
                    else:
                        _set_cell_style(table.cell(row_i, c_idx), "-\n(-)", bg_color=COLOR_YELLOW, font_size=8)
                    c_idx += 1

        # 표 테두리 선 전체 검정색 적용
        for row in table.rows:
            for cell in row.cells:
                _set_cell_border(cell, color="000000")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

